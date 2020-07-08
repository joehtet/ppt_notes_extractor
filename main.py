from pptx import Presentation
from docx import Document
import sys
from docx.shared import Pt

file_name = str(sys.argv[1])
prs = Presentation(file_name)#13800_Introduction to Oncology and Cancer Cell Biology_R5v1.pptx')
document = Document()

style = document.styles['Normal']
font = style.font
font.name = 'Calibri'
font.size = Pt(11)

page_num = 0
start_index = 10 # start index of text after skipping "Narration:" 

def _get_notes_text_list(slide):
    notes_text_list = []

    # Get notes text
    notes_text = ""
    for shape in slide.notes_slide.shapes:
        if shape.text[:10] == "Narration:":
            notes_text = shape.text
            break

    if notes_text != "":
        # Break up text into list
        notes_text_list = notes_text.strip().split('\n')

    return notes_text_list[1:]

for slide in prs.slides:

    page_num += 1
    notes_text_list = _get_notes_text_list(slide)
    
    if notes_text_list:
        document.add_heading("Slide " + str(page_num) + "\n", level=2)
        for item in notes_text_list:
            paragraph = document.add_paragraph(item , style='ListBullet')
            paragraph.paragraph_format.space_after = Pt(6)

document.save(file_name + ".docx")